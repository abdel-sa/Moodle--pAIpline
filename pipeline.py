import json
import sys
import os
import argparse
import re
from datetime import datetime
from typing import List, Dict, Any
from pathlib import Path
from lxml import etree


def read_chapter_file(filepath: str) -> str:
    """
    Reads the content of a given file based on its extension.
    Supports reading from .txt, .md, .pdf, and .pptx files.
    Returns the extracted text as a single string.
    """
    filepath = Path(filepath)
    
    if not filepath.exists():
        raise FileNotFoundError(f"Datei nicht gefunden: {filepath}")
    
    ext = filepath.suffix.lower()
    
    if ext == ".txt" or ext == ".md":
        return filepath.read_text(encoding="utf-8")
    
    elif ext == ".pdf":
        try:
            import PyPDF2
        except ImportError:
            raise ImportError("PyPDF2 nicht installiert. Installiere: pip install PyPDF2")
        
        text = []
        with open(filepath, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text.append(page.extract_text())
        return "\n".join(text)
    
    elif ext == ".pptx":
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("python-pptx nicht installiert. Installiere: pip install python-pptx")
        
        prs = Presentation(filepath)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    
    else:
        raise ValueError(f"Nicht unterstütztes Format: {ext}")


def call_llm(prompt: str, model: str = "mistral") -> str:
    """
    Calls a local Ollama instance with the specified model and prompt.
    Throws an error if Ollama is not reachable or if the request fails.
    Returns the generated response text from the LLM.
    """
    try:
        import requests
    except ImportError:
        raise ImportError("requests nicht installiert. Installiere: pip install requests")
    
    ollama_url = "http://localhost:11434/api/generate"
    
    try:
        response = requests.post(
            ollama_url,
            json={
                "model": model,
                "prompt": prompt,
                "stream": False,
                "temperature": 0.7
            },
            timeout=300
        )
        
        if response.status_code != 200:
            raise ValueError(f"Ollama antwortet mit Status {response.status_code}")
        
        response_text = response.json().get("response", "").strip()
        return response_text
    
    except requests.exceptions.ConnectionError:
        raise ConnectionError(
            "Ollama läuft nicht!\n\n"
            "Starte Ollama in einem NEUEN Terminal:\n"
            "  ollama serve\n\n"
            "Oder installiere Ollama von: https://ollama.ai"
        )


def plan_command(
    chapter_file: str,
    title: str,
    base_category: str,
    output_file: str,
    model: str = "gpt-4o-mini"
) -> None:
    """
    Executes the first pipeline stage (Planning).
    Reads the chapter text, prompts the LLM to create a structured plan of topics 
    and question types, and saves the resulting JSON plan to the output_file.
    """
    print(f"\n{'='*70}")
    print(f"STAGE 1: PLAN GENERIEREN")
    print(f"{'='*70}\n")
    
    print(f"[1] Lese Kapitel: {chapter_file}")
    chapter_text = read_chapter_file(chapter_file)
    print(f"    ✓ {len(chapter_text)} Zeichen geladen\n")
    
    print(f"[2] Analysiere Kapitel mit Ollama ({model})...")
    print(f"    (Das kann 30-60 Sekunden dauern)...")
    
    planner_prompt = f"""Du bist ein Didaktik-Experte und analysierst ein Kapitel für die Erstellung von Prüfungsfragen.

KAPITEL TITEL: {title}
BASIS-KATEGORIE: {base_category}

KAPITELTEXT:
---
{chapter_text[:3000]}
---

AUFGABE:
Erstelle einen strukturierten Plan für Prüfungsfragen zu diesem Kapitel.

Antworte ONLY mit gültigem JSON (ohne Markdown), im folgenden Format:
{{
  "topics": [
    {{
      "topic": "Thema",
      "qtype": "oumultiresponse",
      "num_questions": 3,
      "category_path": "{base_category}/Unterkategorie"
    }}
  ]
}}

FRAGETYPEN:
- oumultiresponse: Multiple Choice mit mehreren richtigen Antworten
- coderunner: Programmieraufgaben
- gapfill: Lückentext
- matching: Zuordnungsaufgaben
- numerical: Numerische Antworten

WICHTIG:
- Mindestens 3, maximal 10 Topics
- topic: Beschreibung des Unterthemas (Deutsch, 3-8 Worte)
- qtype: Einer der obigen Typen
- num_questions: 2-5 pro Topic
- NUR gültiges JSON, keine Markdown!"""
    
    try:
        response = call_llm(planner_prompt, model=model)
        
        if "```json" in response:
            response = response.split("```json")[1].split("```")[0]
        elif "```" in response:
            response = response.split("```")[1].split("```")[0]
        
        plan_data = json.loads(response.strip())
        
    except json.JSONDecodeError as e:
        print(f"    ✗ Fehler beim JSON-Parsing: {e}")
        print(f"    Antwort war: {response[:200]}...")
        sys.exit(1)
    
    print(f"    ✓ Plan generiert: {len(plan_data.get('topics', []))} Topics\n")
    
    print(f"[3] Speichere Plan: {output_file}")
    with open(output_file, "w", encoding="utf-8") as f:
        json.dump(plan_data, f, ensure_ascii=False, indent=2)
    
    print(f"    ✓ Plan gespeichert\n")
    
    print(f"{'='*70}")
    print(f"PLAN SUMMARY:")
    print(f"{'='*70}")
    for i, topic_data in enumerate(plan_data.get('topics', []), 1):
        print(f"  [{i}] {topic_data['topic']}")
        print(f"      Typ: {topic_data['qtype']} | {topic_data['num_questions']} Fragen")
        print(f"      Kategorie: {topic_data['category_path']}")
    print()


def generate_command(
    plan_file: str,
    output_file: str,
    interactive: bool = False,
    model: str = "gpt-4o-mini"
) -> None:
    """
    Executes the second pipeline stage (Generation).
    Loads a previously generated plan, prompts the LLM for individual questions per topic,
    and packages them into a final Moodle XML file for easy importing.
    Supports an interactive mode to review and regenerate questions manually.
    """
    print(f"\n{'='*70}")
    print(f"STAGE 2: FRAGEN GENERIEREN")
    print(f"{'='*70}\n")
    
    print(f"[1] Lade Plan: {plan_file}")
    if not os.path.exists(plan_file):
        print(f"    ✗ Datei nicht gefunden: {plan_file}")
        sys.exit(1)
    
    with open(plan_file, "r", encoding="utf-8") as f:
        plan_data = json.load(f)
    
    print(f"    ✓ Plan geladen: {len(plan_data.get('topics', []))} Topics\n")
    
    print(f"[2] Initialisiere Moodle-XML...")
    xml_root = init_moodle_xml()
    print(f"    ✓ XML-Root erstellt\n")
    
    topics = plan_data.get('topics', [])
    total_questions = 0
    question_counter = 0
    
    for topic_idx, topic_data in enumerate(topics, 1):
        topic = topic_data['topic']
        qtype = topic_data['qtype']
        num_questions = topic_data['num_questions']
        category_path = topic_data['category_path']
        
        print(f"{'─'*70}")
        print(f"Topic [{topic_idx}/{len(topics)}]: {topic}")
        print(f"  Typ: {qtype} | {num_questions} Fragen | Kategorie: {category_path}")
        print(f"{'─'*70}")
        
        questions = generate_questions_for_topic(
            topic=topic,
            qtype=qtype,
            num_questions=num_questions,
            model=model
        )
        
        if not questions:
            print(f"  ✗ Keine Fragen generiert. Überspringe Topic.\n")
            continue
        
        if interactive:
            print(f"\n  Generierte Fragen:")
            for i, q in enumerate(questions, 1):
                print(f"    [{i}] {q.get('name', 'Untitled')}")
            
            while True:
                choice = input("\n  Akzeptieren [a], neu generieren [r], überspringen [s]? ").strip().lower()
                if choice == 'a':
                    break
                elif choice == 'r':
                    questions = generate_questions_for_topic(
                        topic=topic,
                        qtype=qtype,
                        num_questions=num_questions,
                        model=model
                    )
                    for i, q in enumerate(questions, 1):
                        print(f"    [{i}] {q.get('name', 'Untitled')}")
                elif choice == 's':
                    questions = []
                    break
                else:
                    print("  ? Ungültige Eingabe. Versuche nochmal.")
        
        for q in questions:
            question_counter += 1
            add_question_from_json(xml_root, q, qtype, category_path, question_counter)
        
        total_questions += len(questions)
        print(f"  ✓ {len(questions)} Frage(n) hinzugefügt\n")
    
    print(f"[3] Speichere XML: {output_file}")
    save_moodle_xml(xml_root, output_file)
    print(f"    ✓ {total_questions} Fragen exportiert\n")
    
    print(f"{'='*70}")
    print(f"ERFOLG! XML generiert: {output_file}")
    print(f"Fragen importierbar in Moodle → Fragenbank → Import")
    print(f"{'='*70}\n")


def generate_questions_for_topic(
    topic: str,
    qtype: str,
    num_questions: int,
    model: str = "mistral"
) -> List[Dict[str, Any]]:
    """
    Dispatches the LLM call format depending on the required Moodle question type
    (e.g., oumultiresponse, coderunner, gapfill, etc.). 
    Parses the returned JSON and falls back to a default set of questions if errors occur.
    """
    
    if qtype == "oumultiresponse":
        prompt = _prompt_oumultiresponse(topic, num_questions)
    elif qtype == "coderunner":
        prompt = _prompt_coderunner(topic, num_questions)
    elif qtype == "gapfill":
        prompt = _prompt_gapfill(topic, num_questions)
    elif qtype == "matching":
        prompt = _prompt_matching(topic, num_questions)
    elif qtype == "numerical":
        prompt = _prompt_numerical(topic, num_questions)
    else:
        print(f"  ✗ Unbekannter Fragetyp: {qtype}")
        return []
    
    try:
        response = call_llm(prompt, model=model)
        
        if "```json" in response:
            response = response.split("```json")[1].split("```")[0]
        elif "```" in response:
            response = response.split("```")[1].split("```")[0]
        
        response = response.strip()
        
        import re
        response = re.sub(r',(\s*[}\]])', r'\1', response)
        response = response.replace('\n', ' ')
        response = response.replace('  ', ' ')
        
        data = json.loads(response)
        questions = data.get("questions", [])
        
        if not questions:
            questions = _generate_fallback_questions(topic, qtype, num_questions)
        
        return questions
    
    except json.JSONDecodeError as e:
        print(f"  ⚠ JSON-Parsing fehlgeschlagen: {str(e)[:50]}...")
        print(f"    Nutze Fallback-Fragen stattdessen.")
        return _generate_fallback_questions(topic, qtype, num_questions)
    
    except Exception as e:
        print(f"  ✗ Fehler beim Generieren: {e}")
        return _generate_fallback_questions(topic, qtype, num_questions)


def _generate_fallback_questions(topic: str, qtype: str, num_questions: int) -> List[Dict[str, Any]]:
    questions = []
    
    if qtype == "oumultiresponse":
        fallback_qa = {
            "Programmiersprachen": [
                {
                    "name": "Programmiersprachen: Funktionen",
                    "questiontext": "Welche der folgenden Programmiersprachen unterstützen Funktionen?",
                    "answers": [
                        {"text": "Python unterstützt Funktionen.", "correct": True},
                        {"text": "Java unterstützt Funktionen.", "correct": True},
                        {"text": "HTML unterstützt Funktionen.", "correct": False}
                    ]
                }
            ],
            "Python-Datentypen": [
                {
                    "name": "Datentypen: Primitive",
                    "questiontext": "Welche sind primitive Datentypen in Python?",
                    "answers": [
                        {"text": "int (Ganzzahl)", "correct": True},
                        {"text": "str (String)", "correct": True},
                        {"text": "liste (Liste)", "correct": False}
                    ]
                }
            ]
        }
        
        for key in fallback_qa:
            if key.lower() in topic.lower():
                for q in fallback_qa[key][:num_questions]:
                    questions.append(q)
                return questions
        
        for i in range(1, min(num_questions + 1, 3)):
            questions.append({
                "name": f"{topic}: Konzept {i}",
                "questiontext": f"Welche Aussagen zu {topic} sind korrekt?",
                "answers": [
                    {"text": f"Aussage {chr(65+i)} ist richtig.", "correct": True},
                    {"text": f"Aussage {chr(66+i)} ist richtig.", "correct": True},
                    {"text": f"Aussage {chr(67+i)} ist falsch.", "correct": False}
                ]
            })
    
    elif qtype == "matching":
        questions.append({
            "name": f"{topic}: Zuordnung",
            "questiontext": f"Ordne die Begriffe den Definitionen zu:",
            "pairs": [
                {"left": "Konzept 1", "right": "Definition für Konzept 1"},
                {"left": "Konzept 2", "right": "Definition für Konzept 2"}
            ]
        })
    
    elif qtype in ["coderunner", "gapfill"]:
        questions.append({
            "name": f"{topic}: Praktische Aufgabe",
            "questiontext": f"Bearbeite folgende Aufgabe zu {topic}."
        })
    
    elif qtype == "numerical":
        questions.append({
            "name": f"{topic}: Berechnung",
            "questiontext": f"Berechne einen Wert zum Thema {topic}:",
            "correct_answer": 42,
            "tolerance": 0.1
        })
    
    return questions


def _prompt_oumultiresponse(topic: str, num_questions: int) -> str:
    return f"""Generiere genau {num_questions} Multiple-Choice-Fragen zum Thema: {topic}

{{
  "questions": [
    {{
      "name": "{topic}: Eigenschaft 1",
      "questiontext": "Welche der folgenden Aussagen zu {topic} sind korrekt?",
      "answers": [
        {{"text": "Aussage A ist richtig.", "correct": true}},
        {{"text": "Aussage B ist falsch.", "correct": false}}
      ]
    }}
  ]
}}"""


def _prompt_coderunner(topic: str, num_questions: int) -> str:
    return f"""Generiere {num_questions} Programmieraufgaben zum Thema: {topic}

{{
  "questions": [
    {{
      "name": "CodeRunner: {topic} - Aufgabe 1",
      "questiontext": "Schreibe eine Funktion..."
    }}
  ]
}}"""


def _prompt_gapfill(topic: str, num_questions: int) -> str:
    return f"""Generiere {num_questions} Lückentexte zum Thema: {topic}

{{
  "questions": [
    {{
      "name": "Gapfill: {topic} - Text 1",
      "questiontext": "Ein Lückentext..."
    }}
  ]
}}"""


def _prompt_matching(topic: str, num_questions: int) -> str:
    return f"""Generiere {num_questions} Zuordnungsaufgaben zum Thema: {topic}

{{
  "questions": [
    {{
      "name": "Matching: {topic} - Aufgabe 1",
      "questiontext": "Ordne zu:",
      "pairs": [
        {{"left": "A", "right": "Definition A"}}
      ]
    }}
  ]
}}"""


def _prompt_numerical(topic: str, num_questions: int) -> str:
    return f"""Generiere {num_questions} numerische Fragen zum Thema: {topic}

{{
  "questions": [
    {{
      "name": "Numerisch: {topic} - Frage 1",
      "correct_answer": 42,
      "tolerance": 0.01
    }}
  ]
}}"""


def init_moodle_xml():
    return etree.Element("quiz")


def add_question_from_json(
    xml_root,
    question_data: Dict[str, Any],
    qtype: str,
    category_path: str,
    question_id: int
) -> None:
    if qtype == "oumultiresponse":
        add_mc_question(xml_root, question_data, category_path, question_id)
    elif qtype in ("coderunner", "gapfill"):
        add_shortanswer_question(xml_root, question_data, category_path, question_id)
    elif qtype == "matching":
        add_matching_question(xml_root, question_data, category_path, question_id)
    elif qtype == "numerical":
        add_numerical_question(xml_root, question_data, category_path, question_id)



def _init_question(xml_root, moodle_type: str, data: dict, default_name: str, category_path: str):
    """Builds question element with name and questiontext. Returns (question, category_fn).
    Caller adds type-specific children, then calls category_fn() to close."""
    q = etree.SubElement(xml_root, "question")
    q.set("type", moodle_type)
    etree.SubElement(etree.SubElement(q, "name"), "text").text = str(data.get("name", default_name))
    qt = etree.SubElement(q, "questiontext")
    qt.set("format", "html")
    etree.SubElement(qt, "text").text = str(data.get("questiontext", ""))
    return q


def _append_category(question, category_path: str) -> None:
    etree.SubElement(etree.SubElement(question, "category"), "text").text = str(category_path)


def add_mc_question(xml_root, question_data: Dict[str, Any], category_path: str, question_id: int) -> None:
    question = _init_question(xml_root, "multichoice", question_data, f"Frage {question_id}", category_path)
    
    shuffleanswers = etree.SubElement(question, "shuffleanswers")
    shuffleanswers.text = "1"
    
    answernumbering = etree.SubElement(question, "answernumbering")
    answernumbering.text = "abc"
    
    single = etree.SubElement(question, "single")
    single.text = "false"
    
    for idx, answer_data in enumerate(question_data.get("answers", []), 1):
        answer = etree.SubElement(question, "answer")
        fraction = "100" if answer_data.get("correct") else "0"
        answer.set("fraction", fraction)
        answer.set("format", "html")
        
        answer_text = etree.SubElement(answer, "text")
        answer_text.text = str(answer_data.get("text", ""))
        
        feedback = etree.SubElement(answer, "feedback")
        feedback.set("format", "html")
        etree.SubElement(feedback, "text").text = ""

    _append_category(question, category_path)


def add_shortanswer_question(xml_root, question_data: Dict[str, Any], category_path: str, question_id: int) -> None:
    question = _init_question(xml_root, "shortanswer", question_data, f"Frage {question_id}", category_path)
    answer = etree.SubElement(question, "answer")
    answer.set("fraction", "100")
    answer.set("format", "html")
    etree.SubElement(answer, "text").text = "Siehe Fragebeschreibung"
    _append_category(question, category_path)


def add_matching_question(xml_root, question_data: Dict[str, Any], category_path: str, question_id: int) -> None:
    question = _init_question(xml_root, "matching", question_data, f"Matching {question_id}", category_path)
    etree.SubElement(question, "shuffleanswers").text = "1"

    for pair in question_data.get("pairs", []):
        sub = etree.SubElement(question, "subquestion")
        sub.set("format", "html")
        etree.SubElement(sub, "text").text = str(pair.get("left", ""))
        etree.SubElement(etree.SubElement(sub, "answer"), "text").text = str(pair.get("right", ""))

    _append_category(question, category_path)


def add_numerical_question(xml_root, question_data: Dict[str, Any], category_path: str, question_id: int) -> None:
    question = _init_question(xml_root, "numerical", question_data, f"Numerisch {question_id}", category_path)
    answer = etree.SubElement(question, "answer")
    answer.set("fraction", "100")
    answer.set("format", "html")
    etree.SubElement(answer, "text").text = str(question_data.get("correct_answer", "0"))
    etree.SubElement(answer, "tolerance").text = str(question_data.get("tolerance", "0.01"))
    _append_category(question, category_path)


def save_moodle_xml(xml_root, output_file: str) -> None:
    tree = etree.ElementTree(xml_root)
    tree.write(
        output_file,
        encoding="utf-8",
        xml_declaration=True,
        standalone=True,
        pretty_print=True
    )


def main():
    """
    Main entry point for the CLI tool. 
    Parses arguments for the 'plan' and 'generate' subcommands and routes them accordingly.
    """
    parser = argparse.ArgumentParser(description="Moodle-Fragenpipeline")
    subparsers = parser.add_subparsers(dest="command", help="Befehl")
    
    plan_parser = subparsers.add_parser("plan")
    plan_parser.add_argument("--chapter", required=True)
    plan_parser.add_argument("--title", required=True)
    plan_parser.add_argument("--base-category", required=True)
    plan_parser.add_argument("--out", required=True)
    plan_parser.add_argument("--model", default="mistral")
    
    gen_parser = subparsers.add_parser("generate")
    gen_parser.add_argument("--plan", required=True)
    gen_parser.add_argument("--out", required=True)
    gen_parser.add_argument("--interactive", action="store_true")
    gen_parser.add_argument("--model", default="mistral")
    
    args = parser.parse_args()
    
    if not args.command:
        parser.print_help()
        sys.exit(0)
    
    try:
        if args.command == "plan":
            plan_command(
                chapter_file=args.chapter,
                title=args.title,
                base_category=args.base_category,
                output_file=args.out,
                model=args.model
            )
        
        elif args.command == "generate":
            output_file = args.out
            if "{timestamp}" in output_file:
                timestamp = datetime.now().strftime("%Y%m%d-%H%M")
                output_file = output_file.format(timestamp=timestamp)
            
            generate_command(
                plan_file=args.plan,
                output_file=output_file,
                interactive=args.interactive,
                model=args.model
            )
    
    except KeyboardInterrupt:
        print("\n\nAbgebrochen durch Benutzer.")
        sys.exit(0)
    
    except Exception as e:
        print(f"\n✗ Fehler: {e}")
        sys.exit(1)


if __name__ == "__main__":
    main()
